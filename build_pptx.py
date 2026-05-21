"""Build the UniLipi slide deck (mirror of the project webpage)."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

ROOT = Path(__file__).resolve().parent
IMG = ROOT / "static" / "images"
LOGOS = IMG / "logos"
OUT = ROOT / "UniLipi.pptx"

# ----- palette (from style.css) -----
ACCENT     = RGBColor(0x25, 0x63, 0xEB)
ACCENT_DK  = RGBColor(0x1E, 0x40, 0xAF)
INK        = RGBColor(0x1F, 0x29, 0x33)
MUTED      = RGBColor(0x52, 0x60, 0x6D)
RULE       = RGBColor(0xE4, 0xE7, 0xEB)
BG_SOFT    = RGBColor(0xF8, 0xFA, 0xFC)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
NORTHCEN   = RGBColor(0x4A, 0x86, 0xE8)
SOUTH      = RGBColor(0xE0, 0x66, 0x66)
HIMALAYAN  = RGBColor(0x8E, 0x7C, 0xC3)
WESTERN    = RGBColor(0x6A, 0xA8, 0x4F)
EASTERN    = RGBColor(0xF6, 0xB2, 0x6B)
OURS_BG    = RGBColor(0xE0, 0xF2, 0xFE)
OURS_INK   = RGBColor(0x0C, 0x4A, 0x6E)

REGION_COLORS = [NORTHCEN, WESTERN, HIMALAYAN, EASTERN, SOUTH]

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
blank = prs.slide_layouts[6]


# --------- helpers ---------
def set_bg(slide, rgb):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb
    bg.shadow.inherit = False
    return bg


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri", italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    r.font.color.rgb = color
    return tb


def add_rich(slide, x, y, w, h, runs, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.15):
    """runs: list of (text, dict(size,bold,color,italic,font))"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    for i, (text, props) in enumerate(runs):
        r = p.add_run()
        r.text = text
        r.font.size = Pt(props.get("size", 18))
        r.font.bold = props.get("bold", False)
        r.font.italic = props.get("italic", False)
        r.font.name = props.get("font", "Calibri")
        r.font.color.rgb = props.get("color", INK)
    return tb


def add_pill(slide, x, y, w, h, text, *, fill=RGBColor(0xDB, 0xEA, 0xFE),
             color=ACCENT, size=11, bold=True):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.adjustments[0] = 0.5
    s.line.fill.background()
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    tf = s.text_frame
    tf.margin_left = Emu(60000); tf.margin_right = Emu(60000)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Calibri"
    return s


def add_rect(slide, x, y, w, h, fill, *, line=None, radius=None):
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape, x, y, w, h)
    if radius is not None:
        s.adjustments[0] = radius
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(0.75)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    return s


def section_title(slide, text, accent=ACCENT):
    """Vertical accent bar + section title at top of slide."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.45),
                                 Inches(0.08), Inches(0.55))
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    add_text(slide, Inches(0.85), Inches(0.4), Inches(12), Inches(0.65),
             text, size=30, bold=True, color=INK,
             anchor=MSO_ANCHOR.MIDDLE, font="Cambria")


def footer(slide, page_num, total, label):
    add_text(slide, Inches(0.6), Inches(7.05), Inches(8), Inches(0.35),
             "UniLipi  ·  ICDAR 2026 Submission", size=10, color=MUTED)
    add_text(slide, Inches(10.5), Inches(7.05), Inches(2.6), Inches(0.35),
             f"{label}   {page_num} / {total}", size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def fit_image(slide, path, cx, cy, max_w, max_h, *, border=True):
    """Place image centered at (cx,cy) preserving aspect ratio within max_w x max_h."""
    with Image.open(path) as im:
        iw, ih = im.size
    ar = iw / ih
    target_ar = max_w / max_h
    if ar > target_ar:
        w = max_w
        h = int(max_w / ar)
    else:
        h = max_h
        w = int(max_h * ar)
    x = int(cx - w / 2)
    y = int(cy - h / 2)
    if border:
        # subtle white card behind the image
        pad = Inches(0.08)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      x - pad, y - pad, w + 2 * pad, h + 2 * pad)
        card.adjustments[0] = 0.04
        card.fill.solid(); card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RULE; card.line.width = Pt(0.75)
    return slide.shapes.add_picture(str(path), x, y, width=w, height=h)


# =====================================================================
# Slide 1 — Title / hero
# =====================================================================
s = prs.slides.add_slide(blank)
set_bg(s, WHITE)

# soft gradient-ish band at top
band = add_rect(s, 0, 0, SW, Inches(2.6), BG_SOFT)
# colored thin stripe
stripe = add_rect(s, 0, Inches(2.58), SW, Inches(0.04), ACCENT)

# venue tag
add_pill(s, Inches(5.55), Inches(0.55), Inches(2.25), Inches(0.34),
         "SUBMITTED TO ICDAR 2026", size=11)

# title
add_text(s, Inches(0.8), Inches(1.05), Inches(11.75), Inches(1.4),
         "UniLipi: A Unified Multi-Script OCR\nfor Historical Indic Manuscripts",
         size=36, bold=True, color=INK, align=PP_ALIGN.CENTER, font="Cambria")

# subtitle (rich)
add_rich(s, Inches(0.8), Inches(2.45), Inches(11.75), Inches(0.45),
         [("A single OCR foundation model trained jointly across ",
           {"size": 17, "color": MUTED}),
          ("13 Indic scripts", {"size": 17, "color": ACCENT, "bold": True}),
          (".", {"size": 17, "color": MUTED})],
         align=PP_ALIGN.CENTER)

# authors
add_text(s, Inches(0.8), Inches(3.0), Inches(11.75), Inches(0.4),
         "Tathagata Ghosh   ·   Sai Madhusudan Gunda   ·   "
         "Simran Singh Sandral   ·   Ravi Kiran Sarvadevabhatla",
         size=14, color=INK, align=PP_ALIGN.CENTER)

# affiliation
add_text(s, Inches(0.8), Inches(3.4), Inches(11.75), Inches(0.35),
         "International Institute of Information Technology, Hyderabad",
         size=12, italic=True, color=MUTED, align=PP_ALIGN.CENTER)

# stats row
stat_data = [
    ("13",    "Indic scripts",     NORTHCEN),
    ("32.5M", "Synthetic lines",   WESTERN),
    ("6.9%",  "Overall CER",       SOUTH),
    ("85M",   "Parameters",        HIMALAYAN),
]
sw = Inches(2.2); sh = Inches(1.05); gap = Inches(0.25)
total_w = sw * 4 + gap * 3
sx = (SW - total_w) // 2
sy = Inches(4.05)
for i, (val, lbl, col) in enumerate(stat_data):
    x = sx + i * (sw + gap)
    card = add_rect(s, x, sy, sw, sh, WHITE, line=RULE, radius=0.12)
    add_text(s, x, sy + Inches(0.12), sw, Inches(0.5), val,
             size=26, bold=True, color=col, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE, font="Cambria")
    add_text(s, x, sy + Inches(0.62), sw, Inches(0.4), lbl.upper(),
             size=9, bold=True, color=MUTED, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)

# logos row
logos = [
    (LOGOS / "icdar2026.png", Inches(1.0)),
    (LOGOS / "iiith.jpg",     Inches(0.9)),
    (LOGOS / "lekhya.png",    Inches(1.0)),
]
total_logo_w = sum(w for _, w in logos) + Inches(1.2)
lx = (SW - total_logo_w) // 2
ly = Inches(5.6)
for path, w in logos:
    with Image.open(path) as im:
        iw, ih = im.size
    h = int(w * ih / iw)
    if h > Inches(1.1):
        h = Inches(1.1); w = int(h * iw / ih)
    s.shapes.add_picture(str(path), lx, ly + (Inches(1.1) - h) // 2, width=w, height=h)
    lx += w + Inches(0.6)

# bottom contact
add_text(s, Inches(0.8), Inches(6.95), Inches(11.75), Inches(0.4),
         "tathagata.ghosh@research.iiit.ac.in   ·   sai.gunda@research.iiit.ac.in   "
         "·   ssandral@gitam.in   ·   ravi.kiran@iiit.ac.in",
         size=10, color=MUTED, align=PP_ALIGN.CENTER, font="Consolas")


# =====================================================================
# Slide 2 — Abstract
# =====================================================================
s = prs.slides.add_slide(blank)
set_bg(s, WHITE)
section_title(s, "Abstract", REGION_COLORS[0])

abstract_runs = [
    ("Optical character recognition (OCR) for handwritten Indic manuscripts is essential for "
     "large-scale digitization and computational access to manuscript heritage. Existing approaches, "
     "however, are typically developed for one script at a time and require substantial "
     "script-specific customization — limiting scalability and practical deployment.\n\n",
     {"size": 16, "color": INK}),
    ("We present ", {"size": 16, "color": INK}),
    ("UniLipi", {"size": 16, "color": ACCENT, "bold": True}),
    (", a unified multi-script OCR model for handwritten Indic manuscripts trained jointly across ",
     {"size": 16, "color": INK}),
    ("13 Indic scripts", {"size": 16, "color": INK, "bold": True}),
    (" within a single framework. UniLipi directly handles realistic manuscript conditions: "
     "extreme variation in line geometry, large variation in line length, and partial interruptions "
     "from holes, stains, or pictorial illustrations. To operate effectively under "
     "ultra low-resource conditions, the model leverages ",
     {"size": 16, "color": INK}),
    ("script-aware synthetic manuscript data generation",
     {"size": 16, "color": INK, "bold": True}),
    (", substantially reducing reliance on real annotated data.\n\n",
     {"size": 16, "color": INK}),
    ("Beyond historical manuscripts, UniLipi serves as an effective foundational pretrained model — "
     "its learned representations enable strong OCR performance for contemporary Indic handwriting "
     "and extend to several non-Indic scripts including ",
     {"size": 16, "color": INK}),
    ("Tibetan, Italian, Latin, and Chinese",
     {"size": 16, "color": ACCENT, "bold": True}),
    (". In addition to transcription, UniLipi predicts ",
     {"size": 16, "color": INK}),
    ("script identity", {"size": 16, "color": INK, "bold": True}),
    (" and ", {"size": 16, "color": INK}),
    ("per-line native character counts", {"size": 16, "color": INK, "bold": True}),
    (", supporting practical manuscript cataloging workflows.",
     {"size": 16, "color": INK}),
]
add_rich(s, Inches(0.85), Inches(1.3), Inches(11.6), Inches(5.5),
         abstract_runs, line_spacing=1.3)


# =====================================================================
# Slide 3 — Diversity of Indic Manuscripts
# =====================================================================
s = prs.slides.add_slide(blank)
set_bg(s, BG_SOFT)
section_title(s, "Diversity of Indic Manuscripts", REGION_COLORS[1])

# left column text
text_x, text_y = Inches(0.85), Inches(1.3)
text_w, text_h = Inches(4.4), Inches(5.6)
add_rich(s, text_x, text_y, text_w, Inches(2.4),
         [("Indic manuscripts span centuries, regions, and writing substrates. UniLipi targets ",
           {"size": 14, "color": INK}),
          ("13 scripts", {"size": 14, "color": INK, "bold": True}),
          (" drawn from five regional traditions:",
           {"size": 14, "color": INK})],
         line_spacing=1.3)

regions = [
    ("North / Central Indian",  NORTHCEN, "Devanagari · Sharada · Modi"),
    ("Western Indian",          WESTERN,  "Gurmukhi · Jaini"),
    ("Eastern Indian",          EASTERN,  "Bengali · Odia"),
    ("Southern Indian",         SOUTH,    "Grantha · Telugu · Malayalam · Kannada"),
    ("Himalayan / Buddhist",    HIMALAYAN,"Newar · Siddham"),
]
ry = Inches(2.5)
for name, col, scripts in regions:
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, text_x, ry + Inches(0.07),
                             Inches(0.16), Inches(0.16))
    dot.line.fill.background(); dot.fill.solid(); dot.fill.fore_color.rgb = col
    add_text(s, text_x + Inches(0.3), ry, Inches(4.1), Inches(0.3),
             name, size=12, bold=True, color=col)
    add_text(s, text_x + Inches(0.3), ry + Inches(0.28), Inches(4.1), Inches(0.3),
             scripts, size=11, color=MUTED, italic=True)
    ry += Inches(0.62)

# right column image
fit_image(s, IMG / "real-data.jpg",
          cx=Inches(9.45), cy=Inches(3.85),
          max_w=Inches(7.0), max_h=Inches(5.0))

add_text(s, Inches(5.6), Inches(6.5), Inches(7.5), Inches(0.5),
         "Manuscript lines exhibit substantial diversity in substrate (palm leaf, paper), "
         "layout geometry, stroke thickness, ink, and orthographic structure.",
         size=10, italic=True, color=MUTED, align=PP_ALIGN.CENTER)


# =====================================================================
# Slide 4 — Contributions
# =====================================================================
s = prs.slides.add_slide(blank)
set_bg(s, WHITE)
section_title(s, "Contributions", REGION_COLORS[2])

contribs = [
    ("UniLipi",
     "The first unified foundational multi-script OCR model for handwritten Indic scripts, "
     "trained jointly across 13 distinct Indic scripts within a single framework.",
     NORTHCEN),
    ("Script-aware Synthetic Data",
     "A reference-conditioned synthetic manuscript generation approach that reduces reliance on "
     "large volumes of real annotated handwritten data.",
     WESTERN),
    ("Auxiliary Supervision",
     "An OCR system that, in addition to transcription, predicts script identity and per-line "
     "glyph count — useful for downstream cataloging workflows.",
     HIMALAYAN),
]

cy = Inches(1.45)
for i, (head, body, col) in enumerate(contribs):
    y = cy + Inches(i * 1.75)
    card = add_rect(s, Inches(0.85), y, Inches(11.6), Inches(1.5), WHITE,
                    line=RULE, radius=0.05)
    bar = add_rect(s, Inches(0.85), y, Inches(0.15), Inches(1.5), col)
    add_text(s, Inches(1.2), y + Inches(0.15), Inches(11.1), Inches(0.5),
             head, size=18, bold=True, color=col, font="Cambria")
    add_text(s, Inches(1.2), y + Inches(0.65), Inches(11.1), Inches(0.85),
             body, size=13, color=INK)


# =====================================================================
# Slide 5 — Method
# =====================================================================
s = prs.slides.add_slide(blank)
set_bg(s, BG_SOFT)
section_title(s, "Method — UniLipi Architecture", REGION_COLORS[3])

add_rich(s, Inches(0.85), Inches(1.2), Inches(11.6), Inches(0.7),
         [("UniLipi is a unified multi-task OCR model. Given a manuscript line image, it jointly "
           "predicts a text sequence in the unified ", {"size": 13, "color": INK}),
          ("Roman WX", {"size": 13, "color": ACCENT, "bold": True, "font": "Consolas"}),
          (" representation, a ", {"size": 13, "color": INK}),
          ("script identity label", {"size": 13, "color": INK, "bold": True}),
          (", and a per-line ", {"size": 13, "color": INK}),
          ("glyph count", {"size": 13, "color": INK, "bold": True}),
          (".", {"size": 13, "color": INK})],
         line_spacing=1.25)

# architecture image
fit_image(s, IMG / "model.jpg",
          cx=Inches(4.5), cy=Inches(4.5),
          max_w=Inches(7.4), max_h=Inches(4.5))

# right column — loss
box = add_rect(s, Inches(8.6), Inches(2.3), Inches(4.3), Inches(4.0),
               WHITE, line=RULE, radius=0.05)
add_text(s, Inches(8.8), Inches(2.45), Inches(4.0), Inches(0.4),
         "Joint training objective", size=13, bold=True, color=INK, font="Cambria")
add_text(s, Inches(8.8), Inches(2.95), Inches(4.0), Inches(0.6),
         "L = L_OCR  +  λ_script · L_script\n        +  λ_count · L_count",
         size=12, color=ACCENT, font="Consolas")
add_text(s, Inches(8.8), Inches(4.0), Inches(4.0), Inches(2.2),
         "Multi-task supervision encourages the encoder to learn richer representations that "
         "capture textual content, script identity, and structural length information together — "
         "improving transcription accuracy across diverse manuscripts.",
         size=11, color=INK)


# =====================================================================
# Slide 6 — Synthetic Data Pipeline
# =====================================================================
s = prs.slides.add_slide(blank)
set_bg(s, WHITE)
section_title(s, "Reference-Conditioned Synthetic Data", REGION_COLORS[4])

add_rich(s, Inches(0.85), Inches(1.2), Inches(11.6), Inches(0.7),
         [("Real annotated manuscript lines are extremely scarce for many Indic scripts. We "
           "synthesize training data by lifting ", {"size": 13, "color": INK}),
          ("appearance and geometric priors", {"size": 13, "color": INK, "bold": True}),
          (" from real reference manuscripts — yielding ~", {"size": 13, "color": INK}),
          ("2.5M synthetic lines per script", {"size": 13, "color": ACCENT, "bold": True}),
          (" and ", {"size": 13, "color": INK}),
          ("32.5M total", {"size": 13, "color": ACCENT, "bold": True}),
          (" across 13 scripts.", {"size": 13, "color": INK})],
         line_spacing=1.25)

fit_image(s, IMG / "synthetic.jpg",
          cx=SW // 2, cy=Inches(4.4),
          max_w=Inches(12.4), max_h=Inches(4.4))

# 4-step caption strip
steps = [
    ("1", "Sample Roman WX line and transliterate to target script.", NORTHCEN),
    ("2", "Extract writing mask, font size, ink color from a reference.", WESTERN),
    ("3", "Composite rendered text onto diffusion-inpainted background.", HIMALAYAN),
    ("4", "Output line + WX text + script label + glyph-count.", EASTERN),
]
bw = Inches(2.9); bh = Inches(0.65); gap = Inches(0.13)
bx = (SW - (bw * 4 + gap * 3)) // 2
by = Inches(6.65)
for i, (n, txt, col) in enumerate(steps):
    x = bx + i * (bw + gap)
    add_rect(s, x, by, bw, bh, WHITE, line=col, radius=0.18)
    num = add_rect(s, x + Inches(0.08), by + Inches(0.12), Inches(0.4), Inches(0.4), col,
                   radius=0.5)
    add_text(s, x + Inches(0.08), by + Inches(0.12), Inches(0.4), Inches(0.4),
             n, size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.55), by, bw - Inches(0.6), bh, txt,
             size=9, color=INK, anchor=MSO_ANCHOR.MIDDLE)


# =====================================================================
# Slide 7 — Results (overall) + Ablation
# =====================================================================
s = prs.slides.add_slide(blank)
set_bg(s, BG_SOFT)
section_title(s, "Results — Overall Scores", REGION_COLORS[0])

# left: big-number summary
big_x = Inches(0.85)
big_y = Inches(1.4)
big_w = Inches(5.5)

add_rect(s, big_x, big_y, big_w, Inches(5.5), WHITE, line=RULE, radius=0.05)

add_text(s, big_x + Inches(0.3), big_y + Inches(0.2), big_w - Inches(0.6), Inches(0.5),
         "UniLipi (weighted across 13 scripts)",
         size=14, bold=True, color=INK, font="Cambria")

# CER
add_text(s, big_x + Inches(0.3), big_y + Inches(0.85), big_w - Inches(0.6), Inches(1.0),
         "6.9%", size=54, bold=True, color=ACCENT, font="Cambria")
add_text(s, big_x + Inches(0.3), big_y + Inches(1.85), big_w - Inches(0.6), Inches(0.35),
         "CER  ↓   OCR transcription error",
         size=11, color=MUTED, bold=True)

# Mono CER
add_text(s, big_x + Inches(0.3), big_y + Inches(2.35), big_w - Inches(0.6), Inches(0.6),
         "8.6%", size=28, bold=True, color=SOUTH, font="Cambria")
add_text(s, big_x + Inches(0.3), big_y + Inches(2.95), big_w - Inches(0.6), Inches(0.3),
         "Mono CER  ↓   per-script monolingual baseline",
         size=11, color=MUTED, bold=True)

# Count-MAE
add_text(s, big_x + Inches(0.3), big_y + Inches(3.4), big_w - Inches(0.6), Inches(0.6),
         "1.1", size=28, bold=True, color=WESTERN, font="Cambria")
add_text(s, big_x + Inches(0.3), big_y + Inches(4.0), big_w - Inches(0.6), Inches(0.3),
         "Count-MAE  ↓   per-line glyph count error",
         size=11, color=MUTED, bold=True)

add_text(s, big_x + Inches(0.3), big_y + Inches(4.55), big_w - Inches(0.6), Inches(0.8),
         "Unified multi-script training outperforms dedicated per-script models.",
         size=11, italic=True, color=INK)

# right: ablation table
abl_x = Inches(6.6); abl_y = Inches(1.4); abl_w = Inches(6.3)
add_text(s, abl_x, abl_y, abl_w, Inches(0.4),
         "Ablation Study", size=15, bold=True, color=INK, font="Cambria")
add_text(s, abl_x, abl_y + Inches(0.4), abl_w, Inches(0.35),
         "Single-component swaps vs. final UniLipi design.",
         size=10, italic=True, color=MUTED)

abl_rows = [
    ("Modified Component",       "Default",         "CER ↓",  "MAE ↓", True,  False),
    ("IAST as vocab",            "Roman WX",        "11.0",   "1.7",   False, False),
    ("Dedicated vocab / script", "Roman WX",        "14.1",   "1.9",   False, False),
    ("No script token",          "Auxiliary tokens","7.1",    "1.3",   False, False),
    ("No count token",           "Auxiliary tokens","9.3",    "—",     False, False),
    ("No script & count",        "Auxiliary tokens","9.7",    "—",     False, False),
    ("Encoder — 4 layers",       "6 layers",        "9.3",    "1.5",   False, False),
    ("Encoder — 8 layers",       "6 layers",        "8.1",    "1.3",   False, False),
    ("Cross-Entropy loss",       "CTC loss",        "8.3",    "1.4",   False, False),
    ("UniLipi (Final)",          "Default design",  "6.9",    "1.1",   False, True),
]
col_w = [Inches(2.6), Inches(1.9), Inches(0.9), Inches(0.9)]
row_h = Inches(0.33)
ry = abl_y + Inches(0.85)
for row in abl_rows:
    label, default, cer, mae, is_head, is_ours = row
    bgc = RGBColor(0xF1, 0xF5, 0xF9) if is_head else (OURS_BG if is_ours else WHITE)
    add_rect(s, abl_x, ry, sum(col_w, Emu(0)), row_h, bgc, line=RULE)
    cx = abl_x
    fontc = OURS_INK if is_ours else INK
    for j, txt in enumerate([label, default, cer, mae]):
        align = PP_ALIGN.RIGHT if j >= 2 else PP_ALIGN.LEFT
        add_text(s, cx + Inches(0.1), ry, col_w[j] - Inches(0.15), row_h, txt,
                 size=10, bold=(is_head or is_ours), color=fontc,
                 align=align, anchor=MSO_ANCHOR.MIDDLE)
        cx += col_w[j]
    ry += row_h


# =====================================================================
# Slide 8 — Qualitative Results
# =====================================================================
s = prs.slides.add_slide(blank)
set_bg(s, WHITE)
section_title(s, "Qualitative Results", REGION_COLORS[1])

fit_image(s, IMG / "results.jpg",
          cx=SW // 2, cy=Inches(4.1),
          max_w=Inches(12.4), max_h=Inches(5.0))

add_rich(s, Inches(0.85), Inches(6.7), Inches(11.6), Inches(0.5),
         [("UniLipi transcription on a representative script from each region.  ",
           {"size": 12, "color": INK, "italic": True}),
          ("Red", {"size": 12, "color": SOUTH, "bold": True}),
          (" = deletions   ·   ", {"size": 12, "color": INK, "italic": True}),
          ("Blue", {"size": 12, "color": ACCENT, "bold": True}),
          (" = substitutions / additions.",
           {"size": 12, "color": INK, "italic": True})],
         align=PP_ALIGN.CENTER)


# =====================================================================
# Slide 9 — Heatmap
# =====================================================================
s = prs.slides.add_slide(blank)
set_bg(s, BG_SOFT)
section_title(s, "Count-Token Attention Heatmap", REGION_COLORS[2])

fit_image(s, IMG / "heatmap.jpg",
          cx=Inches(4.5), cy=Inches(4.4),
          max_w=Inches(7.5), max_h=Inches(5.0))

# right caption / takeaway
box_x = Inches(8.7); box_y = Inches(1.45)
add_rect(s, box_x, box_y, Inches(4.2), Inches(5.5), WHITE, line=RULE, radius=0.05)
add_text(s, box_x + Inches(0.3), box_y + Inches(0.2), Inches(3.8), Inches(0.5),
         "What does the count token learn?",
         size=14, bold=True, color=INK, font="Cambria")
add_text(s, box_x + Inches(0.3), box_y + Inches(0.8), Inches(3.8), Inches(4.5),
         "The count token's attention map localizes glyph-level structure along the line. "
         "This serves as:\n\n"
         "•  auxiliary supervision during training, giving the encoder a structural prior;\n\n"
         "•  a lightweight quality-control cue at inference — large gaps between "
         "predicted token count and visual glyph count flag likely OCR errors.",
         size=11, color=INK)


# =====================================================================
# Slide 10 — Thank You / contact
# =====================================================================
s = prs.slides.add_slide(blank)
set_bg(s, WHITE)
band = add_rect(s, 0, 0, SW, Inches(2.5), BG_SOFT)
stripe = add_rect(s, 0, Inches(2.48), SW, Inches(0.04), ACCENT)

add_text(s, Inches(0.6), Inches(0.8), Inches(12.1), Inches(1.0),
         "Thank you",
         size=54, bold=True, color=INK, align=PP_ALIGN.CENTER, font="Cambria")
add_text(s, Inches(0.6), Inches(1.8), Inches(12.1), Inches(0.5),
         "Questions and early-access requests are welcome.",
         size=15, italic=True, color=MUTED, align=PP_ALIGN.CENTER)

# contact card
cx0 = Inches(2.5); cy0 = Inches(3.2); cw = Inches(8.3); chh = Inches(2.2)
add_rect(s, cx0, cy0, cw, chh, WHITE, line=RULE, radius=0.04)
add_pill(s, cx0 + Inches(0.4), cy0 + Inches(0.3), Inches(1.4), Inches(0.34),
         "CONTACT", size=10)
add_text(s, cx0 + Inches(0.4), cy0 + Inches(0.75), cw - Inches(0.8), Inches(0.5),
         "sai.gunda@research.iiit.ac.in",
         size=20, bold=True, color=ACCENT, font="Consolas")
add_text(s, cx0 + Inches(0.4), cy0 + Inches(1.3), cw - Inches(0.8), Inches(0.4),
         "Project page  ·  unilipi (coming soon)",
         size=12, color=MUTED)
add_text(s, cx0 + Inches(0.4), cy0 + Inches(1.65), cw - Inches(0.8), Inches(0.4),
         "Code  ·  github.com/lekhya-ai/UniLipi",
         size=12, color=MUTED, font="Consolas")

# bottom logos
logos2 = [
    (LOGOS / "icdar2026.png", Inches(0.85)),
    (LOGOS / "iiith.jpg",     Inches(0.8)),
    (LOGOS / "lekhya.png",    Inches(0.85)),
]
total_w = sum(w for _, w in logos2) + Inches(1.2)
lx = (SW - total_w) // 2
ly = Inches(5.95)
for path, w in logos2:
    with Image.open(path) as im:
        iw, ih = im.size
    h = int(w * ih / iw)
    if h > Inches(0.95):
        h = Inches(0.95); w = int(h * iw / ih)
    s.shapes.add_picture(str(path), lx, ly + (Inches(0.95) - h) // 2, width=w, height=h)
    lx += w + Inches(0.6)


# add footers to interior slides
total = len(prs.slides)
labels = [
    "Title", "Abstract", "Diversity", "Contributions", "Method",
    "Synthetic", "Results", "Qualitative", "Heatmap", "Thanks"
]
for i, slide in enumerate(prs.slides, start=1):
    if i == 1 or i == total:
        continue
    footer(slide, i, total, labels[i - 1])

prs.save(OUT)
print(f"wrote {OUT}  ({OUT.stat().st_size/1024:.0f} KB, {total} slides)")
